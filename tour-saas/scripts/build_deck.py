#!/usr/bin/env python3
"""Dựng deck PPTX thuyết trình cho dự án tour-saas, rồi LibreOffice convert sang PDF."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Bảng màu thương hiệu ----
TEAL   = RGBColor(0x0F, 0x76, 0x6E)
CYAN   = RGBColor(0x0E, 0x74, 0x90)
INK    = RGBColor(0x1F, 0x29, 0x37)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x04, 0x78, 0x57)
RED    = RGBColor(0xB9, 0x1C, 0x1C)
AMBER  = RGBColor(0xB4, 0x53, 0x09)
CARD   = RGBColor(0xE2, 0xE8, 0xF0)
FONT   = "Liberation Sans"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(tf_para, text, size, color, bold=False, italic=False):
    tf_para.text = text
    for r in tf_para.runs:
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def content_slide(title, eyebrow=None):
    s = prs.slide_layouts and prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 1.15, TEAL)               # thanh tiêu đề
    rect(s, 0, 1.15, 13.333, 0.06, CYAN)
    tf = box(s, 0.6, 0.18, 12.1, 0.85)
    if eyebrow:
        p0 = tf.paragraphs[0]; _set(p0, eyebrow.upper(), 11, RGBColor(0xCF,0xEF,0xEC), bold=True)
        p1 = tf.add_paragraph(); _set(p1, title, 26, WHITE, bold=True)
    else:
        _set(tf.paragraphs[0], title, 28, WHITE, bold=True)
    return s


def bullets(slide, items, left=0.7, top=1.55, width=12.0, size=16, gap=10):
    tf = box(slide, left, top, width, 5.6)
    first = True
    for it in items:
        lvl = it[0] if isinstance(it, tuple) else 0
        txt = it[1] if isinstance(it, tuple) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        mark = "▸ " if lvl == 0 else "– "
        _set(p, mark + txt, size if lvl == 0 else size-2,
             INK if lvl == 0 else MUTED, bold=(lvl == 0 and txt.endswith(":")))
    return tf


# ============ SLIDE 1 — TITLE ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, TEAL)
rect(s, 0, 5.9, 13.333, 1.6, CYAN)
tf = box(s, 0.9, 1.7, 11.5, 3.2)
_set(tf.paragraphs[0], "SIÊU DỰ ÁN · TRAVEL-TECH VIỆT NAM", 14, RGBColor(0xCF,0xEF,0xEC), bold=True)
p = tf.add_paragraph(); _set(p, "Trợ lý Báo giá & Lịch trình Tour bằng AI", 40, WHITE, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(10)
_set(p, "SaaS nội bộ cho agency du lịch Việt Nam tự vận hành", 20, RGBColor(0xE0,0xF2,0xF1))
tf2 = box(s, 0.9, 6.15, 11.5, 1.0)
_set(tf2.paragraphs[0], "Tổng hợp nghiên cứu thị trường (có trích dẫn) + MVP đã chạy được", 13, WHITE)
p = tf2.add_paragraph(); _set(p, "Nhánh: claude/travel-booking-saas-ai-dURHV  ·  Pull Request #1", 12, RGBColor(0xD0,0xEC,0xEA))

# ============ SLIDE 2 — CƠ HỘI ============
s = content_slide("Cơ hội thị trường: thời điểm vàng", "Bối cảnh")
bullets(s, [
  "81% du khách Việt dự định dùng AI cho chuyến đi tới — CAO NHẤT châu Á (Agoda 2026)",
  "Outbound bùng nổ: ~6,7 triệu lượt năm 2025 (+26%), chi trung bình ~36 triệu đ/chuyến",
  "Nội địa: hơn 110 triệu lượt khách năm 2024",
  "OTA ngoại chiếm ~80% thị trường online → cửa thắng là TOOLING B2B cho operator, KHÔNG phải OTA tiêu dùng",
  "Zalo: ~77 triệu MAU — 'đường ray' giao dịch & chốt tour ở Việt Nam",
], size=17, gap=14)

# ============ SLIDE 3 — NỖI ĐAU ============
s = content_slide("Nỗi đau hiện tại của agency", "Vấn đề")
bullets(s, [
  "Sale tốn 1–2 GIỜ để soạn một báo giá thủ công: tra giá khách sạn → gõ Excel → tính markup → format PDF",
  "Rớt khách ngoài giờ; quy trình thủ công, dễ sai số, khó nhân bản",
  "Phần mềm incumbent VN (TravelMaster…) mạnh sổ sách/chiết tính nhưng KHÔNG có AI, KHÔNG Zalo native",
  "Tool AI quốc tế (Mindtrip, Layla…) chỉ là 'máy tạo ý tưởng' cho khách lẻ:",
  (1, "Không có inventory/giá Việt Nam, không báo giá có cost/margin, không brand agency"),
], size=17, gap=14)

# ============ SLIDE 4 — GIẢI PHÁP / WEDGE ============
s = content_slide("Giải pháp — mũi nhọn (wedge)", "Định vị")
tf = box(s, 0.7, 1.5, 12.0, 1.5)
_set(tf.paragraphs[0],
     "Nhập yêu cầu khách (tiếng Việt) → AI dựng lịch trình + chiết tính giá trên KHO NHÀ CUNG CẤP RIÊNG "
     "→ proposal có brand → chốt qua Zalo", 19, TEAL, bold=True)
y = 3.25
for title, desc in [
    ("Dữ liệu giá đã đàm phán riêng", "Kho net rate của agency — thứ các tool toàn cầu không có"),
    ("Tiếng Việt + bối cảnh địa phương", "Hiểu yêu cầu, viết lịch trình tự nhiên cho khách Việt"),
    ("Báo giá có cost/margin + Zalo", "Chiết tính chuẩn nghiệp vụ; chốt trên kênh người Việt dùng"),
]:
    rect(s, 0.7, y, 11.9, 1.05, LIGHT)
    rect(s, 0.7, y, 0.12, 1.05, TEAL)
    tfx = box(s, 1.0, y+0.12, 11.4, 0.85)
    _set(tfx.paragraphs[0], "MOAT · " + title, 15, INK, bold=True)
    p = tfx.add_paragraph(); _set(p, desc, 13, MUTED)
    y += 1.2

# ============ SLIDE 5 — KHOẢNG TRỐNG (TABLE) ============
s = content_slide("Khoảng trống: không ai làm trọn combo", "Đối thủ")
rows = [
    ["Tiêu chí", "ToursMS", "TravelMaster", "AI quốc tế", "DỰ ÁN NÀY"],
    ["AI lịch trình tiếng Việt", "Quảng cáo", "Không", "Không có VN", "Có"],
    ["Kho giá riêng + markup/VAT", "Quảng cáo", "Mạnh", "Không", "Có"],
    ["Tích hợp Zalo", "Quảng cáo", "Không", "Không", "Roadmap"],
    ["Dấu vết / độ tin cậy", "Không có", "Đã kiểm chứng", "Khách lẻ", "Tự kiểm soát"],
]
cols = 5
tbl_shape = s.shapes.add_table(len(rows), cols, Inches(0.7), Inches(1.55),
                               Inches(11.93), Inches(4.0))
table = tbl_shape.table
table.columns[0].width = Inches(3.1)
for c in range(1, cols):
    table.columns[c].width = Inches(2.2075)
for r in range(len(rows)):
    table.rows[r].height = Inches(0.8)
def cell_color(txt):
    if txt in ("Có", "Mạnh", "Đã kiểm chứng"): return GREEN
    if txt in ("Không", "Không có", "Không có VN"): return RED
    if txt in ("Quảng cáo", "Khách lẻ"): return AMBER
    if txt in ("Roadmap", "Tự kiểm soát"): return CYAN
    return INK
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
            _set(para, val, 13, WHITE, bold=True)
        elif c == cols-1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEC,0xFE,0xFF)
            _set(para, val, 12.5, cell_color(val), bold=True)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            if c == 0:
                _set(para, val, 12.5, INK, bold=True)
            else:
                _set(para, val, 12.5, cell_color(val), bold=(val in ("Mạnh","Đã kiểm chứng")))
tf = box(s, 0.7, 6.7, 12.0, 0.6)
_set(tf.paragraphs[0],
     "ToursMS hứa đủ nhưng KHÔNG có dấu vết độc lập nào · TravelMaster mạnh chiết tính nhưng thiếu AI & Zalo",
     12, MUTED, italic=True)

# ============ SLIDE 6 — SẢN PHẨM (LUỒNG) ============
s = content_slide("Sản phẩm — luồng vận hành", "Cách hoạt động")
steps = [
    ("1", "Yêu cầu khách", "Tiếng Việt, qua web / Zalo / tin nhắn"),
    ("2", "AI dựng lịch trình", "Theo ngày + chọn dịch vụ từ kho giá riêng"),
    ("3", "Chiết tính giá", "net → markup → VAT → giá bán → margin"),
    ("4", "Proposal có brand", "Link đẹp gửi khách / xuất PDF"),
]
x = 0.7; w = 2.95; gap = 0.12
for n, t, d in steps:
    rect(s, x, 2.0, w, 3.2, LIGHT)
    rect(s, x, 2.0, w, 0.7, TEAL)
    tfn = box(s, x, 2.05, w, 0.6); tfn.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set(tfn.paragraphs[0], "BƯỚC " + n, 13, WHITE, bold=True)
    tft = box(s, x+0.15, 2.95, w-0.3, 2.1)
    _set(tft.paragraphs[0], t, 16, INK, bold=True)
    p = tft.add_paragraph(); p.space_before = Pt(8); _set(p, d, 13, MUTED)
    x += w + gap
tf = box(s, 0.7, 5.6, 12.0, 1.2)
_set(tf.paragraphs[0], "Nguyên tắc khóa: GIÁ đến từ kho nhà cung cấp — AI chỉ CHỌN dịch vụ, KHÔNG bịa giá.",
     16, TEAL, bold=True)

# ============ SLIDE 7 — KIẾN TRÚC ============
s = content_slide("Kiến trúc kỹ thuật", "Bên dưới")
bullets(s, [
  "LLM (Claude) điều phối + structured JSON output cho lịch trình",
  "Kho giá NCC = 'hard data' (RAG-ready); AI chỉ chọn serviceId có thật → chống bịa giá/hallucinate",
  "Quote engine TÁCH RIÊNG, có unit test — phần tiền phải chuẩn như TravelMaster",
  "Proposal renderer: HTML có brand → in PDF; Fallback theo luật khi không có API key",
  "Stack:",
  (1, "Node.js / TypeScript · Express · @anthropic-ai/sdk (claude-opus-4-8) · pgvector (roadmap)"),
], size=16, gap=13)

# ============ SLIDE 8 — MVP ============
s = content_slide("MVP đã build — chạy được ngay", "Hiện trạng")
bullets(s, [
  "Đã chạy end-to-end, nằm trong Pull Request #1",
  "Kho giá mẫu 6 điểm đến: Đà Nẵng, Hạ Long/Hà Nội, Phú Quốc, Tokyo, Bangkok",
  "AI sinh lịch trình + chiết tính giá + proposal có brand (link đẹp / PDF)",
  "Phần chiết tính giá: 6/6 unit test PASS",
  "Chạy được CẢ KHI không có API key (chế độ fallback theo luật)",
  "Demo Phú Quốc 4 khách: net 24,8tr → tổng 32,19tr (markup 18%, VAT 10%) · lợi nhuận gộp 15,25%",
], size=16, gap=12)

# ============ SLIDE 9 — CHI PHÍ ============
s = content_slide("Chi phí vận hành — rất nhẹ", "Kinh tế")
bullets(s, [
  "~3,7¢ / báo giá (Haiku) → ~10–20¢ (Opus) — chọn model theo nhu cầu chất lượng vs chi phí",
  "Embed cả kho giá ~1 USD một lần; vector DB pgvector gần như miễn phí",
  "Tổng vận hành ở volume thấp: ~20–150 USD/tháng — hoàn toàn trong tầm solo/team nhỏ",
  "Prompt caching giảm mạnh chi phí phần input lặp lại (kho giá + hướng dẫn)",
], size=17, gap=15)

# ============ SLIDE 10 — BUILD VS BUY ============
s = content_slide("Build vs Buy — khuyến nghị", "Chiến lược")
bullets(s, [
  "KHÔNG buy mù: ToursMS/TourPRO chưa có review, khách hàng, hay app store nào kiểm chứng được",
  "Buy có kiểm chứng: dùng thử ToursMS (14 ngày) + TravelMaster (trial) để benchmark",
  "BUILD phần lõi giá trị (moat): AI tiếng Việt + kho giá riêng + Zalo",
  "Học chuẩn chiết tính của TravelMaster (net→markup→VAT); chưa cần API vé/khách sạn live ở giai đoạn đầu",
], size=17, gap=14)

# ============ SLIDE 11 — ROADMAP ============
s = content_slide("Lộ trình phát triển", "Roadmap")
phases = [
    ("Giai đoạn 1 — XONG", "MVP: báo giá + lịch trình AI + proposal có brand", GREEN),
    ("Giai đoạn 2", "Tích hợp Zalo OA · Quản lý giá theo mùa / loại phòng", CYAN),
    ("Giai đoạn 3", "Giá vé/KS live: Duffel (có VNA), RateHawk/TBO, Travelfusion (Vietjet)", TEAL),
    ("Giai đoạn 4", "Postgres + pgvector (RAG kho lớn) · Xuất PDF server-side", INK),
]
y = 1.7
for ph, desc, col in phases:
    rect(s, 0.7, y, 0.12, 1.1, col)
    rect(s, 0.95, y, 11.65, 1.1, LIGHT)
    tfx = box(s, 1.2, y+0.13, 11.2, 0.9)
    _set(tfx.paragraphs[0], ph, 15, col, bold=True)
    p = tfx.add_paragraph(); _set(p, desc, 14, INK)
    y += 1.27

# ============ SLIDE 12 — KẾT LUẬN ============
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, TEAL)
rect(s, 0, 0, 13.333, 0.18, CYAN)
tf = box(s, 0.9, 0.9, 11.5, 1.0)
_set(tf.paragraphs[0], "Kết luận & bước tiếp theo", 32, WHITE, bold=True)
items = [
  "Khoảng trống thị trường có thật — và Việt Nam là nơi khát AI du lịch nhất châu Á",
  "MVP đã chứng minh luồng lõi; chi phí vận hành rất nhẹ (~20–150 USD/tháng)",
  "Moat = dữ liệu giá riêng + tiếng Việt + Zalo — thứ đối thủ không có",
  "Tiếp theo: benchmark ToursMS/TravelMaster → build Zalo + giá theo mùa",
]
tfb = box(s, 0.95, 2.2, 11.4, 3.6)
first = True
for it in items:
    p = tfb.paragraphs[0] if first else tfb.add_paragraph(); first = False
    p.space_after = Pt(16)
    _set(p, "▸  " + it, 18, WHITE)
rect(s, 0, 6.6, 13.333, 0.9, CYAN)
tff = box(s, 0.9, 6.72, 11.5, 0.7)
_set(tff.paragraphs[0], "Mã nguồn: nhánh claude/travel-booking-saas-ai-dURHV · Pull Request #1", 13, WHITE)

import os
out = os.path.join(os.path.dirname(__file__), "..", "docs", "tour-saas-presentation.pptx")
out = os.path.abspath(out)
prs.save(out)
print("Saved", out, "·", len(prs.slides._sldIdLst), "slides")
